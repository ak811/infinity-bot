# cogs/engagement/quiz_maker/file_loader.py
from __future__ import annotations

import csv
from io import BytesIO, StringIO
from pathlib import Path
from typing import Iterable

try:
    from PyPDF2 import PdfReader
except Exception:  # pragma: no cover
    PdfReader = None  # type: ignore[assignment]

# Optional support for Word documents
try:
    import docx  # python-docx
except Exception:  # pragma: no cover
    docx = None  # type: ignore[assignment]

# Optional support for PowerPoint
try:
    from pptx import Presentation  # python-pptx
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore[assignment]

# Optional support for Excel
try:
    import openpyxl
except Exception:  # pragma: no cover
    openpyxl = None  # type: ignore[assignment]


def _decode_bytes(data: bytes) -> str:
    """Decode bytes to text using utf-8 with fallback."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1", errors="ignore")


def _csv_to_text(data: bytes) -> str:
    """Convert CSV bytes into a rough text representation."""
    decoded = _decode_bytes(data)
    reader = csv.reader(StringIO(decoded))
    lines: Iterable[str] = (", ".join(row) for row in reader)
    return "\n".join(lines)


def _pdf_to_text(data: bytes) -> str:
    """Extract text from a PDF file using PyPDF2."""
    if PdfReader is None:
        raise RuntimeError(
            "PyPDF2 is required for PDF support. Install with 'pip install PyPDF2'."
        )

    bio = BytesIO(data)
    reader = PdfReader(bio)
    parts: list[str] = []
    for page in reader.pages:
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        if txt:
            parts.append(txt)
    return "\n".join(parts)


def _docx_to_text(data: bytes) -> str:
    """Extract text from a .docx file using python-docx."""
    if docx is None:
        raise RuntimeError(
            "python-docx is required for .docx support. Install with 'pip install python-docx'."
        )

    bio = BytesIO(data)
    document = docx.Document(bio)
    parts: list[str] = []

    # Paragraphs
    for para in document.paragraphs:
        text = (para.text or "").strip()
        if text:
            parts.append(text)

    # Tables
    for table in document.tables:
        for row in table.rows:
            cell_texts = [
                (cell.text or "").strip()
                for cell in row.cells
                if (cell.text or "").strip()
            ]
            if cell_texts:
                parts.append(" | ".join(cell_texts))

    return "\n".join(parts)


def _pptx_to_text(data: bytes) -> str:
    """Extract text from a .pptx file using python-pptx."""
    if Presentation is None:
        raise RuntimeError(
            "python-pptx is required for .pptx support. Install with 'pip install python-pptx'."
        )

    bio = BytesIO(data)
    prs = Presentation(bio)
    parts: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            text = text.strip()
            if text:
                parts.append(text)

    return "\n".join(parts)


def _xlsx_to_text(data: bytes) -> str:
    """Extract text from an Excel workbook using openpyxl."""
    if openpyxl is None:
        raise RuntimeError(
            "openpyxl is required for Excel support. Install with 'pip install openpyxl'."
        )

    bio = BytesIO(data)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    parts: list[str] = []

    for ws in wb.worksheets:
        parts.append(f"Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            values = [str(v) for v in row if v is not None]
            if values:
                parts.append(", ".join(values))

    return "\n".join(parts)


def extract_text_from_file(filename: str, data: bytes) -> str:
    """
    Best effort extraction of text from a file based on extension.

    Tries to handle:
    - Plain text and code: txt, md, log, json, yaml, yml, ini, cfg, toml, py, js, ts, html, htm, css
    - CSV
    - PDF
    - DOCX
    - PPTX
    - XLSX and related formats

    Falls back to plain decoding for unknown types.
    """
    suffix = Path(filename).suffix.lower()

    text_like = {
        ".txt",
        ".md",
        ".log",
        ".json",
        ".yaml",
        ".yml",
        ".ini",
        ".cfg",
        ".toml",
        ".py",
        ".js",
        ".ts",
        ".html",
        ".htm",
        ".css",
    }

    if suffix in text_like:
        return _decode_bytes(data)

    if suffix == ".csv":
        return _csv_to_text(data)

    if suffix == ".pdf":
        return _pdf_to_text(data)

    if suffix == ".docx":
        return _docx_to_text(data)

    if suffix == ".pptx":
        return _pptx_to_text(data)

    if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        return _xlsx_to_text(data)

    # Generic fallback for anything else
    return _decode_bytes(data)
